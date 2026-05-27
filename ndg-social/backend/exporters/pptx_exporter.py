"""
pptx_exporter — 이미지 방식
  각 슬라이드를 PNG로 스크린샷하여 PPT 슬라이드 전체에 삽입
"""

import asyncio
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

EXPORTS_DIR = Path(__file__).parent.parent.parent / "exports"
EXPORTS_DIR.mkdir(exist_ok=True)

PX_TO_EMU   = 6350
SLIDE_W_EMU = 1920 * PX_TO_EMU
SLIDE_H_EMU = 1080 * PX_TO_EMU


async def _render_and_capture(print_url: str) -> list[bytes]:
    from playwright.async_api import async_playwright
    import logging
    log = logging.getLogger(__name__)

    screenshots = []
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page(viewport={"width": 1920, "height": 1080}, device_scale_factor=3)

        try:
            await page.goto(print_url, wait_until="networkidle", timeout=45000)
        except Exception as e:
            log.warning(f"[PPTX] goto 경고: {e}")

        try:
            await page.wait_for_function("document.title === 'slides-ready'", timeout=20000)
        except Exception:
            log.warning("[PPTX] 'slides-ready' 대기 시간 초과 — 강제 진행")
        await page.wait_for_timeout(1000)

        slide_els = await page.query_selector_all(".slide-print-page")
        log.warning(f"[PPTX] slide_count={len(slide_els)}")

        if not slide_els:
            await browser.close()
            raise RuntimeError("슬라이드 요소 없음")

        for slide_el in slide_els:
            screenshots.append(await slide_el.screenshot(type="png"))

        await browser.close()

    return screenshots


def _lock_picture(pic_shape) -> None:
    """배경 이미지를 선택 불가(noSelect)로 잠금 — 클릭이 텍스트박스로 통과됨"""
    from lxml import etree
    from pptx.oxml.ns import qn

    nvPicPr  = pic_shape._element.find(qn('p:nvPicPr'))
    if nvPicPr is None:
        return
    cNvPicPr = nvPicPr.find(qn('p:cNvPicPr'))
    if cNvPicPr is None:
        return
    picLocks = cNvPicPr.find(qn('a:picLocks'))
    if picLocks is None:
        picLocks = etree.SubElement(cNvPicPr, qn('a:picLocks'))
    picLocks.set('noSelect', '1')
    picLocks.set('noChangeAspect', '1')


def _build_pptx(screenshots: list[bytes]) -> BytesIO:
    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank = prs.slide_layouts[6]

    for png in screenshots:
        slide = prs.slides.add_slide(blank)
        img_stream = BytesIO(png)
        slide.shapes.add_picture(
            img_stream,
            left=Emu(0), top=Emu(0),
            width=Emu(SLIDE_W_EMU), height=Emu(SLIDE_H_EMU),
        )

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


async def export_pptx(print_url: str, filename: str) -> Path:
    output_path = EXPORTS_DIR / filename
    try:
        screenshots = await _render_and_capture(print_url)
    except ImportError:
        raise RuntimeError(
            "Playwright가 설치되지 않았습니다.\n"
            "pip install playwright && playwright install chromium"
        )
    if not screenshots:
        raise RuntimeError("슬라이드를 캡처할 수 없습니다.")

    pptx_buf = _build_pptx(screenshots)
    with open(output_path, "wb") as f:
        f.write(pptx_buf.read())
    return output_path


def export_pptx_sync(print_url: str, filename: str) -> Path:
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor() as pool:
                future = pool.submit(asyncio.run, export_pptx(print_url, filename))
                return future.result()
        else:
            return loop.run_until_complete(export_pptx(print_url, filename))
    except RuntimeError:
        return asyncio.run(export_pptx(print_url, filename))
